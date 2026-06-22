"""Parse uploaded documents into normalized plain text.

One parser per type, dispatched by file extension. Parser libs are imported
lazily so the package imports without the optional `[parsers]` extra installed;
only actually parsing a given type requires its library.
"""

from __future__ import annotations

import html
import io
import os
import re
import tempfile
from collections.abc import Callable
from pathlib import Path

from cue.ingestion.errors import DocumentParseError, UnsupportedDocTypeError
from cue.ingestion.models import DocType

# Strip C0/C1 control chars but keep tab (\x09) and newline (\x0a).
_CONTROL_CHARS = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]")
_HTML_BLOCKS = re.compile(r"(?is)<(script|style)\b.*?</\1>")
_HTML_TAGS = re.compile(r"(?s)<[^>]+>")

_EXTENSION_TO_TYPE: dict[str, DocType] = {
    ".pdf": DocType.pdf,
    ".docx": DocType.docx,
    ".txt": DocType.txt,
    ".pptx": DocType.pptx,
    ".epub": DocType.epub,
}


def detect_doc_type(filename: str) -> DocType:
    """Map a filename's extension to a supported `DocType`, else raise."""
    ext = Path(filename).suffix.lower()
    doc_type = _EXTENSION_TO_TYPE.get(ext)
    if doc_type is None:
        supported = ", ".join(sorted(_EXTENSION_TO_TYPE))
        raise UnsupportedDocTypeError(
            f"Unsupported file type '{ext or filename}'. Supported: {supported}."
        )
    return doc_type


def normalize_text(raw: str) -> str:
    """Collapse whitespace and drop control chars; preserve line structure."""
    text = _CONTROL_CHARS.sub("", raw)
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.splitlines()]
    text = "\n".join(lines)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _strip_html(markup: str) -> str:
    markup = _HTML_BLOCKS.sub(" ", markup)
    markup = _HTML_TAGS.sub(" ", markup)
    return html.unescape(markup)


def parse_txt(data: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def parse_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    try:
        reader = PdfReader(io.BytesIO(data))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as exc:  # noqa: BLE001 - normalize any library failure
        raise DocumentParseError(f"Failed to read PDF: {exc}") from exc


def parse_docx(data: bytes) -> str:
    from docx import Document

    try:
        document = Document(io.BytesIO(data))
        return "\n".join(p.text for p in document.paragraphs)
    except Exception as exc:  # noqa: BLE001
        raise DocumentParseError(f"Failed to read DOCX: {exc}") from exc


def parse_pptx(data: bytes) -> str:
    from pptx import Presentation

    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        raise DocumentParseError(f"Failed to read PPTX: {exc}") from exc

    chunks: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text:
                chunks.append(shape.text_frame.text)
    return "\n".join(chunks)


def parse_epub(data: bytes) -> str:
    import ebooklib
    from ebooklib import epub

    # ebooklib reads from a path, so stage the bytes in a temp file first.
    try:
        with tempfile.TemporaryDirectory() as tmp:
            path = os.path.join(tmp, "book.epub")
            with open(path, "wb") as handle:
                handle.write(data)
            book = epub.read_epub(path)
    except Exception as exc:  # noqa: BLE001
        raise DocumentParseError(f"Failed to read EPUB: {exc}") from exc

    parts = [
        _strip_html(item.get_content().decode("utf-8", errors="replace"))
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT)
    ]
    return "\n".join(parts)


_PARSERS: dict[DocType, Callable[[bytes], str]] = {
    DocType.pdf: parse_pdf,
    DocType.docx: parse_docx,
    DocType.txt: parse_txt,
    DocType.pptx: parse_pptx,
    DocType.epub: parse_epub,
}


def parse(filename: str, data: bytes) -> tuple[DocType, str]:
    """Detect type, parse, and normalize. Returns `(doc_type, normalized_text)`."""
    doc_type = detect_doc_type(filename)
    raw = _PARSERS[doc_type](data)
    return doc_type, normalize_text(raw)
