"""Generate tiny in-memory sample documents for ingestion tests.

Everything is built at runtime so no binary fixtures live in the repo.
"""

from __future__ import annotations

import io
import os
import tempfile


def make_txt_bytes(text: str) -> bytes:
    return text.encode("utf-8")


def make_docx_bytes(text: str) -> bytes:
    from docx import Document

    document = Document()
    document.add_paragraph(text)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def make_pptx_bytes(text: str) -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Title Only
    slide.shapes.title.text = text
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def make_epub_bytes(text: str) -> bytes:
    from ebooklib import epub

    book = epub.EpubBook()
    book.set_identifier("cue-test")
    book.set_title("Cue Test Book")
    book.set_language("en")

    chapter = epub.EpubHtml(title="Chapter 1", file_name="chap1.xhtml", lang="en")
    chapter.content = f"<html><body><h1>Chapter 1</h1><p>{text}</p></body></html>"
    book.add_item(chapter)
    book.toc = (chapter,)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav", chapter]

    with tempfile.TemporaryDirectory() as tmp:
        path = os.path.join(tmp, "book.epub")
        epub.write_epub(path, book)
        with open(path, "rb") as handle:
            return handle.read()


def make_pdf_bytes(text: str) -> bytes:
    """Build a minimal single-page PDF with one line of extractable text.

    `text` must be Latin-1 encodable (it goes straight into a PDF string).
    The cross-reference offsets are computed so pypdf reads it cleanly.
    """
    content = b"BT /F1 24 Tf 72 700 Td (" + text.encode("latin-1") + b") Tj ET"
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
        b"<< /Length " + str(len(content)).encode() + b" >>\nstream\n" + content + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]

    out = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for number, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += str(number).encode() + b" 0 obj\n" + body + b"\nendobj\n"

    xref_pos = len(out)
    size = len(objects) + 1
    out += b"xref\n0 " + str(size).encode() + b"\n"
    out += b"0000000000 65535 f \n"
    for offset in offsets:
        out += f"{offset:010d} 00000 n \n".encode()
    out += b"trailer\n<< /Size " + str(size).encode() + b" /Root 1 0 R >>\n"
    out += b"startxref\n" + str(xref_pos).encode() + b"\n%%EOF"
    return bytes(out)
