"""
Document processing service.
Handles uploading, parsing, chunking, and embedding documents
for the RAG pipeline.
"""

import os
import uuid
import logging
from pathlib import Path
from typing import BinaryIO

from app.core.config import get_settings
from app.models.schemas import DocumentType, DocumentStatus, DocumentResponse

logger = logging.getLogger(__name__)

# Supported file extensions mapped to document types
EXTENSION_MAP: dict[str, DocumentType] = {
    ".pdf": DocumentType.PDF,
    ".docx": DocumentType.DOCX,
    ".txt": DocumentType.TXT,
    ".epub": DocumentType.EPUB,
    ".ppt": DocumentType.PPT,
    ".pptx": DocumentType.PPTX,
}

# Upload directory for temporary file storage
UPLOAD_DIR = Path("./uploads")
UPLOAD_DIR.mkdir(exist_ok=True)


class DocumentService:
    """Service for processing and managing user-uploaded documents."""

    def __init__(self):
        self._documents: dict[str, DocumentResponse] = {}  # in-memory store for MVP

    def get_document_type(self, filename: str) -> DocumentType:
        """Determine document type from file extension."""
        ext = Path(filename).suffix.lower()
        if ext not in EXTENSION_MAP:
            raise ValueError(
                f"Unsupported file type: {ext}. "
                f"Supported: {', '.join(EXTENSION_MAP.keys())}"
            )
        return EXTENSION_MAP[ext]

    async def save_upload(self, filename: str, file: BinaryIO) -> str:
        """Save an uploaded file to disk and return the file path."""
        doc_id = f"doc_{uuid.uuid4().hex[:12]}"
        safe_filename = f"{doc_id}_{filename}"
        file_path = UPLOAD_DIR / safe_filename

        content = await file.read()
        with open(file_path, "wb") as f:
            f.write(content)

        logger.info(f"Saved upload: {safe_filename} ({len(content)} bytes)")
        return str(file_path), doc_id, len(content)

    async def extract_text(self, file_path: str, doc_type: DocumentType) -> str:
        """Extract text content from a document based on its type."""
        if doc_type == DocumentType.PDF:
            return await self._extract_pdf(file_path)
        elif doc_type == DocumentType.DOCX:
            return await self._extract_docx(file_path)
        elif doc_type == DocumentType.TXT:
            return await self._extract_txt(file_path)
        elif doc_type == DocumentType.EPUB:
            return await self._extract_epub(file_path)
        elif doc_type in (DocumentType.PPT, DocumentType.PPTX):
            return await self._extract_pptx(file_path)
        else:
            raise ValueError(f"Unsupported document type: {doc_type}")

    async def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF using PyPDF2."""
        try:
            from PyPDF2 import PdfReader

            reader = PdfReader(file_path)
            text_parts = []
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Failed to extract PDF: {e}")
            raise

    async def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX using python-docx."""
        try:
            from docx import Document

            doc = Document(file_path)
            return "\n\n".join(para.text for para in doc.paragraphs if para.text.strip())
        except Exception as e:
            logger.error(f"Failed to extract DOCX: {e}")
            raise

    async def _extract_txt(self, file_path: str) -> str:
        """Read plain text files."""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    async def _extract_epub(self, file_path: str) -> str:
        """Extract text from EPUB using ebooklib + BeautifulSoup."""
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup

            book = epub.read_epub(file_path)
            text_parts = []
            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                soup = BeautifulSoup(item.get_content(), "html.parser")
                text = soup.get_text(separator="\n")
                if text.strip():
                    text_parts.append(text.strip())
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Failed to extract EPUB: {e}")
            raise

    async def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PPTX using python-pptx."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = [f"[Slide {slide_num}]"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                text_parts.append("\n".join(slide_texts))
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Failed to extract PPTX: {e}")
            raise

    def register_document(self, doc_id: str, filename: str, doc_type: DocumentType,
                          size_bytes: int, chunk_count: int = 0,
                          status: DocumentStatus = DocumentStatus.PROCESSING) -> DocumentResponse:
        """Register a document in the in-memory store."""
        doc = DocumentResponse(
            id=doc_id,
            filename=filename,
            type=doc_type,
            size_bytes=size_bytes,
            chunk_count=chunk_count,
            status=status,
        )
        self._documents[doc_id] = doc
        return doc

    def update_status(self, doc_id: str, status: DocumentStatus, chunk_count: int = 0):
        """Update a document's processing status."""
        if doc_id in self._documents:
            self._documents[doc_id].status = status
            if chunk_count:
                self._documents[doc_id].chunk_count = chunk_count

    def get_document(self, doc_id: str) -> DocumentResponse | None:
        """Retrieve a document by ID."""
        return self._documents.get(doc_id)

    def list_documents(self) -> list[DocumentResponse]:
        """List all documents."""
        return list(self._documents.values())

    def delete_document(self, doc_id: str) -> bool:
        """Delete a document from the store."""
        if doc_id in self._documents:
            del self._documents[doc_id]
            return True
        return False


# Singleton instance
document_service = DocumentService()
